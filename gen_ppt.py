from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 定义颜色
PRIMARY_COLOR = RgbColor(0x1a, 0x52, 0x76)
WHITE = RgbColor(0xFF, 0xFF, 0xFF)
TEXT_COLOR = RgbColor(0x33, 0x33, 0x33)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RgbColor(0x5d, 0xad, 0xe2)
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_lines):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY_COLOR
    shape.line.fill.background()
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, line in enumerate(content_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(20)
        p.font.color.rgb = TEXT_COLOR
        p.space_before = Pt(10)
    
    return slide

# 生成PPT页面
add_title_slide(prs, "皮诊智核", "基于RAG与Agent的皮肤病变智能诊断系统\n\nAI赋能皮肤疾病全流程智能诊疗")

add_content_slide(prs, "目录", [
    "1. 项目背景与痛点",
    "2. 项目概述", 
    "3. 技术架构总览",
    "4. 核心技术亮点",
    "5. 系统功能展示",
    "6. 模型性能评估",
    "7. 创新点总结",
    "8. 应用场景",
    "9. 项目成果与展望"
])

add_content_slide(prs, "1. 项目背景与痛点", [
    "皮肤疾病诊疗现状的三大痛点：",
    "",
    "医疗资源不均 — 皮肤科医生集中于大城市，基层和偏远地区诊断能力薄弱",
    "诊断主观性强 — 肉眼观察依赖经验，易漏诊误诊", 
    "患者就医困难 — 挂号难、排队久、费用高，小问题也需跑医院",
    "",
    "数据支撑：我国皮肤科医生约3万人，面对近14亿人口，医患比严重失衡"
])

add_content_slide(prs, "2. 项目概述", [
    "皮诊智核是什么？",
    "",
    "一款融合 深度学习图像识别 + RAG知识增强 + 多Agent智能协作 的皮肤疾病AI辅助诊疗系统",
    "",
    "核心能力：",
    "精准图像识别 — 7类常见皮肤病，准确率>90%",
    "智能症状分析 — 自然语言理解",
    "医学知识检索 — RAG增强，可溯源",
    "个性化治疗建议 — 多Agent协作决策"
])

add_content_slide(prs, "3. 技术架构总览", [
    "四层架构设计：",
    "",
    "用户交互层：React前端 — 现代化UI，响应式设计",
    "",
    "智能协作层：多Agent系统 — 症状分析Agent + 图像诊断Agent + 治疗建议Agent",
    "              多Agent协作管理器 — ReAct推理 + 自我验证",
    "",
    "模型推理层：深度学习模型 — EfficientNet-B3 / ResNet50 / CustomSkinNet / YOLO",
    "",
    "知识增强层：RAG系统 — Chroma向量数据库 + 多源检索 + 引用标注"
])

add_content_slide(prs, "4. 核心技术亮点 ① — 深度学习模型", [
    "多模型融合策略：",
    "",
    "EfficientNet-B3 — 精度-效率平衡最优，作为主分类器",
    "ResNet50 — 经典稳定，迁移学习友好，用于辅助验证",
    "CustomSkinNet — 自研轻量级网络，支持边缘设备部署",
    "YOLO — 实时目标检测，实现病灶定位",
    "",
    "数据基础：HAM10000权威数据集，7类皮肤病变，10015张医学图像"
])

add_content_slide(prs, "4. 核心技术亮点 ② — RAG知识增强", [
    "为什么需要RAG？",
    "",
    "传统LLM：可能幻觉，答案不可信",
    "RAG增强：检索真实医学文献，答案可溯源",
    "",
    "技术流程：",
    "多源检索 — 从医学知识库召回相关文档",
    "智能重排 — 按相关性排序优化结果",
    "多重过滤 — 去重、去噪、保质量",
    "引用标注 — 每个答案标注参考来源",
    "来源追溯 — 支持查证和进一步阅读"
])

add_content_slide(prs, "4. 核心技术亮点 ③ — 多Agent智能协作", [
    "模仿真实医生会诊流程：",
    "",
    "用户输入 → 症状分析Agent → 图像诊断Agent → 治疗建议Agent",
    "                ↓",
    "           自我验证模块（置信度检查 / 紧急情况识别 / 不确定性提醒）",
    "                ↓",
    "            输出报告",
    "",
    "ReAct推理：每一步决策都有思考过程可追溯"
])

add_content_slide(prs, "5. 系统功能展示", [
    "三大核心功能模块：",
    "",
    "图像诊断 — 上传皮肤照片，AI识别病变类型，3秒出结果，初筛效率高",
    "",
    "症状咨询 — 描述症状，AI智能分析，自然交互，无需医学术语",
    "",
    "综合报告 — 融合图像+症状，生成诊断建议，个性化、可溯源、可打印",
    "",
    "特色功能：紧急情况自动识别并提醒就医 / 诊断置信度透明展示 / 支持多轮对话追问"
])

add_content_slide(prs, "5. 系统功能展示 — 前端界面", [
    "技术栈：React 18 + Vite + Ant Design + ECharts",
    "",
    "界面特点：",
    "现代化UI设计，医疗级专业感",
    "响应式布局，支持PC/平板",
    "数据可视化（诊断统计、置信度图表）",
    "报告导出（PDF/图片）"
])

add_content_slide(prs, "6. 模型性能评估", [
    "分类性能指标：",
    "",
    "准确率 (Accuracy) — > 90%",
    "精确率 (Precision) — 待补充",
    "召回率 (Recall) — 待补充",
    "F1-Score — 待补充",
    "Top-5 准确率 — 待补充",
    "",
    "支持7类皮肤病变分类：akiec、bcc、bkl、df、mel、nv、vasc",
    "",
    "混淆矩阵已生成，支持详细性能分析"
])

add_content_slide(prs, "7. 创新点总结", [
    "四大核心创新：",
    "",
    "技术融合创新 — 首次将RAG与多Agent架构引入皮肤疾病诊疗",
    "",
    "可解释AI — ReAct推理+自我验证，诊断过程透明可追溯",
    "",
    "知识增强 — 医学文献检索+引用标注，答案专业可信",
    "",
    "全流程覆盖 — 从症状输入到治疗建议，端到端智能服务"
])

add_content_slide(prs, "8. 应用场景", [
    "三大应用场景：",
    "",
    "基层医疗辅助 — 乡镇卫生院、社区卫生服务中心，辅助全科医生进行皮肤病初筛",
    "",
    "个人健康管理 — 家庭自测，小病早发现，减少不必要的医院奔波",
    "",
    "互联网医疗 — 在线问诊平台集成，提升远程诊疗效率"
])

add_content_slide(prs, "9. 竞品对比", [
    "皮诊智核 vs 传统AI诊断：",
    "",
    "技术架构：单一CNN分类 → 深度学习+RAG+多Agent",
    "",
    "知识来源：训练数据固定 → 动态检索医学文献",
    "",
    "可解释性：黑盒输出 → 推理过程透明",
    "",
    "交互方式：仅图像上传 → 图文结合+多轮对话",
    "",
    "诊断报告：简单标签 → 综合报告+治疗建议"
])

add_content_slide(prs, "10. 项目成果与展望", [
    "已完成：",
    "7类皮肤病变识别模型",
    "RAG知识增强系统",
    "多Agent协作框架",
    "React前端界面",
    "完整诊断流程打通",
    "",
    "未来规划：",
    "扩展至更多皮肤病类型（目标30+类）",
    "移动端APP开发",
    "与医院HIS系统对接",
    "多模态融合（图像+症状+病史）"
])

add_content_slide(prs, "总结", [
    "核心价值：",
    "",
    "让优质皮肤诊疗服务触手可及",
    "",
    "降低就医门槛",
    "提升诊断效率",
    "赋能基层医疗",
    "守护全民皮肤健康"
])

add_title_slide(prs, "感谢聆听！", "Q & A\n\n皮诊智核 — 基于RAG与Agent的皮肤病变智能诊断系统")

# 保存PPT
output_path = r"E:\py项目\Skin diseases\比赛PPT_皮诊智核.pptx"
prs.save(output_path)
print(f"PPT已生成: {output_path}")
